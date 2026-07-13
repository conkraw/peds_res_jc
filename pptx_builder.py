"""PowerPoint export logic for Journal Club Builder."""

from __future__ import annotations

import math 
from io import BytesIO
from typing import Any, Dict, Iterable, List

import qrcode

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from slide_schema import SLIDES
from feedback_config import (
    FEEDBACK_INSTRUCTION,
    REDCAP_DISPLAY_URL,
    REDCAP_QR_URL,
    THANK_YOU_MESSAGE,
    THANK_YOU_TITLE,
)




# Simple, conservative styling. Change these values if you want a Penn State-like theme.
COLOR_DARK = RGBColor(35, 35, 35)
COLOR_MID = RGBColor(95, 95, 95)
COLOR_LIGHT_GRAY = RGBColor(242, 242, 242)
COLOR_HEADER = RGBColor(70, 70, 70)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_ACCENT = RGBColor(30, 80, 130)
COLOR_ACCENT_LIGHT = RGBColor(226, 236, 246)
COLOR_WARNING_LIGHT = RGBColor(250, 238, 218)

SLIDE_W = 13.333
SLIDE_H = 7.5
CUSTOM_SLIDES_KEY = "_custom_slides"
DEFAULT_CUSTOM_SLIDE_INSERT_AFTER = "final_bottom_line"

# Small callout labels should feel visually connected to their content without
# sitting on top of it. Content boxes keep their fixed coordinates, while the
# label is lifted slightly and filled boxes get consistent inner padding.
SMALL_LABEL_HEIGHT = 0.22
SMALL_LABEL_LIFT = 0.07
CALLOUT_MARGIN = 0.11
INLINE_LABEL_HEIGHT = 0.22


def _safe_text(value: Any) -> str:
    if value is None:
        return ""
    return str(value)


def _lines(value: Any) -> List[str]:
    return [line.strip() for line in _safe_text(value).splitlines() if line.strip()]

def _core_slide_ids() -> List[str]:
    return [slide["id"] for slide in SLIDES]


def _normalize_custom_slide_position(value: Any) -> str:
    candidate = str(value or "").strip()
    if candidate in _core_slide_ids():
        return candidate
    return DEFAULT_CUSTOM_SLIDE_INSERT_AFTER


def _custom_slide_insert_after(custom_slide: Dict[str, Any]) -> str:
    return _normalize_custom_slide_position(custom_slide.get("insert_after") or custom_slide.get("after_slide_id"))


def _core_slide_label(slide_id: str) -> str:
    for slide in SLIDES:
        if slide.get("id") == slide_id:
            return str(slide.get("export_title") or slide.get("label") or slide_id)
    return "Final Bottom Line"


def _ordered_custom_slides_after(deck: Dict[str, Any], slide_id: str) -> List[Dict[str, Any]]:
    custom_slides = deck.get(CUSTOM_SLIDES_KEY, [])
    if not isinstance(custom_slides, list):
        return []
    return [
        custom_slide
        for custom_slide in custom_slides
        if isinstance(custom_slide, dict) and _custom_slide_insert_after(custom_slide) == slide_id
    ]


def estimate_textbox_height(
    text: Any,
    width_inches: float,
    font_size: int,
    min_height: float = 0.45,
    max_height: float = 2.2,
    padding: float = 0.25,
) -> float:
    """
    Estimate textbox height in inches based on text length, box width, and font size.
    This is not pixel-perfect, but it works well enough for PowerPoint layout.
    """

    clean_text = _safe_text(text).strip()
    if not clean_text:
        return min_height

    # Bigger font = fewer characters per line.
    # Wider box = more characters per line.
    chars_per_line = max(20, int(width_inches * (105 / font_size)))

    wrapped_lines = 0
    for paragraph in clean_text.splitlines():
        paragraph = paragraph.strip()
        if not paragraph:
            wrapped_lines += 1
        else:
            wrapped_lines += math.ceil(len(paragraph) / chars_per_line)

    line_height = font_size * 0.018
    estimated = padding + wrapped_lines * line_height

    return max(min_height, min(max_height, estimated))

def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: Any,
    font_size: int = 20,
    bold: bool = False,
    color: RGBColor = COLOR_DARK,
    align=PP_ALIGN.LEFT,
    fill: RGBColor | None = None,
    margin: float = 0.08,
    breathing_room: bool | None = None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    # Filled callout boxes receive a little extra internal padding and shrink
    # text only when needed. This keeps the boxes fixed in place while giving
    # the text breathing room, so manual edits usually involve text—not moving
    # shapes around the slide.
    use_breathing_room = (fill is not None) if breathing_room is None else breathing_room
    effective_margin = max(float(margin), CALLOUT_MARGIN) if use_breathing_room else float(margin)
    tf.margin_left = Inches(effective_margin)
    tf.margin_right = Inches(effective_margin)
    tf.margin_top = Inches(effective_margin)
    tf.margin_bottom = Inches(effective_margin)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if use_breathing_room else MSO_ANCHOR.TOP
    if use_breathing_room:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _safe_text(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color

    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = fill

    return shape


def add_title(slide, title: str, subtitle: str | None = None):
    add_textbox(slide, 0.55, 0.22, 12.2, 0.5, title, font_size=30, bold=True, color=COLOR_DARK)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(0.82), Inches(12.1), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.color.rgb = COLOR_ACCENT
    if subtitle:
        add_textbox(slide, 0.65, 0.9, 12.0, 0.4, subtitle, font_size=16, color=COLOR_MID)


def add_bullets(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    items: Iterable[str],
    font_size: int = 18,
    bullet: bool = True,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)

    usable_items = [str(item).strip() for item in items if str(item).strip()]
    if not usable_items:
        usable_items = [""]

    for i, item in enumerate(usable_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_DARK
        p.level = 0
        # python-pptx doesn't expose bullet toggling consistently across versions;
        # prefixing is predictable and remains editable.
        if bullet and not item.startswith(("•", "-", "A.", "B.", "C.", "D.", "1.", "2.", "3.", "4.", "5.")):
            p.text = f"• {item}"

    return shape


def add_section_label(slide, x: float, y: float, w: float, label: str, fill: RGBColor = COLOR_ACCENT_LIGHT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.38))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.bold = True
    r.font.size = Pt(18)
    r.font.color.rgb = COLOR_DARK
    return shape


def add_small_label(slide, x: float, y: float, w: float, label: str, color: RGBColor = COLOR_ACCENT):
    """Add a visible, compact label above a callout box.

    The earlier label was just tiny text on a white background, so it could look
    like it was missing after export. This version draws a small editable
    rounded-rectangle "pill" label. It is lifted slightly above the supplied
    position, leaving a narrow visual gap before the content box. The label
    width is based on the text, not the full callout width.
    """
    label_text = _safe_text(label).strip()
    if not label_text:
        return None

    # Approximate width in inches: enough for the label, but never wider than
    # the space the caller gave us. This keeps labels compact and consistent.
    label_w = min(float(w), max(1.35, min(4.8, 0.12 * len(label_text) + 0.35)))
    label_h = SMALL_LABEL_HEIGHT

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(max(0.0, y - SMALL_LABEL_LIFT)),
        Inches(label_w),
        Inches(label_h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT_LIGHT
    shape.line.color.rgb = COLOR_ACCENT_LIGHT

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label_text
    r.font.bold = True
    r.font.size = Pt(10.5)
    r.font.color.rgb = color

    # Mark these shapes so they can be moved to the top of the z-order after
    # all adjoining content boxes have been added to the slide.
    shape.name = f"JC_LABEL__{label_text}"
    return shape


def add_inline_label(
    slide,
    x: float,
    y: float,
    w: float,
    label: str,
    color: RGBColor = COLOR_ACCENT,
):
    """Add a clean text-only label above content that has no filled box.

    Attached pill labels work well when they visually belong to a shaded
    callout box. They look detached when the content below is plain text.
    This text-only treatment keeps the hierarchy without creating a floating
    rounded rectangle.
    """
    label_text = _safe_text(label).strip()
    if not label_text:
        return None

    shape = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(INLINE_LABEL_HEIGHT),
    )
    shape.name = f"JC_INLINE_LABEL__{label_text}"

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = label_text
    r.font.bold = True
    r.font.size = Pt(11)
    r.font.color.rgb = color
    return shape


def bring_small_labels_to_front(slide) -> None:
    """Move all small callout labels above overlapping content boxes.

    Several layouts intentionally place a label like a tab along the top edge
    of the content box beneath it. PowerPoint uses creation order as z-order,
    so a later content box can otherwise cover part of the label. Re-appending
    only the marked label elements makes those labels the topmost editable
    shapes without changing any coordinates.
    """
    labels = [
        shape
        for shape in slide.shapes
        if str(getattr(shape, "name", "")).startswith("JC_LABEL__")
    ]

    sp_tree = slide.shapes._spTree
    for shape in labels:
        element = shape._element
        parent = element.getparent()
        if parent is None:
            continue
        parent.remove(element)

        # Keep an optional p:extLst element last, as required by the schema.
        insert_at = len(sp_tree)
        if len(sp_tree) and sp_tree[-1].tag.endswith("}extLst"):
            insert_at -= 1
        sp_tree.insert(insert_at, element)


def add_footer(slide, text: str = "Journal Club Builder"):
    add_textbox(slide, 0.6, 7.08, 12.0, 0.24, text, font_size=8, color=COLOR_MID, align=PP_ALIGN.RIGHT)
    bring_small_labels_to_front(slide)


def add_results_table(slide, rows_data: List[Dict[str, Any]], x=0.55, y=1.35, w=12.25, h=3.55):
    default_columns = ["Outcome", "88% threshold", "92% threshold", "Difference", "Interpretation"]

    # Use the column order coming from the Streamlit data editor. This allows
    # Slide 4 table columns to be customized and still export as editable PPT text.
    columns = list(rows_data[0].keys()) if rows_data else default_columns
    if not columns:
        columns = default_columns
    columns = columns[:8]  # keep the slide readable

    cleaned_rows = []
    for row in rows_data or []:
        if any(_safe_text(row.get(col, "")).strip() for col in columns):
            cleaned_rows.append({col: _safe_text(row.get(col, "")) for col in columns})
    cleaned_rows = cleaned_rows[:5]
    if not cleaned_rows:
        cleaned_rows = [{col: "" for col in columns}]

    shape = slide.shapes.add_table(
        len(cleaned_rows) + 1,
        len(columns),
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    table = shape.table

    if len(columns) == 5:
        widths = [2.65, 1.55, 1.55, 1.8, 4.7]
    else:
        widths = [w / len(columns)] * len(columns)
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)

    for c, col in enumerate(columns):
        cell = table.cell(0, c)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER
        cell.margin_left = Inches(0.04)
        cell.margin_right = Inches(0.04)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(10 if len(columns) > 5 else 11)
                run.font.color.rgb = COLOR_WHITE

    for r, row in enumerate(cleaned_rows, start=1):
        for c, col in enumerate(columns):
            cell = table.cell(r, c)
            cell.text = _safe_text(row.get(col, ""))
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT if c == 0 or c == len(columns) - 1 else PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(9.5 if len(columns) > 5 else 10.5)
                    run.font.color.rgb = COLOR_DARK

    return shape


def add_big_number_card(slide, big_number: str, caption: str):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.1), Inches(1.55), Inches(9.2), Inches(2.55))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_ACCENT_LIGHT
    card.line.color.rgb = COLOR_ACCENT

    add_textbox(slide, 2.35, 1.82, 8.7, 0.8, big_number, font_size=44, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, 2.55, 2.82, 8.3, 0.8, caption, font_size=20, bold=False, align=PP_ALIGN.CENTER)


def add_simple_bar_chart(slide, title: str, label1: str, value1: float, label2: str, value2: float, units: str):
    add_textbox(slide, 0.8, 1.25, 11.8, 0.35, title, font_size=18, bold=True, align=PP_ALIGN.CENTER)
    max_value = max(float(value1 or 0), float(value2 or 0), 1.0)
    chart_x = 2.3
    chart_y = 1.9
    chart_w = 8.8
    bar_h = 0.52
    gap = 0.75

    def draw_bar(row_idx: int, label: str, value: float):
        y = chart_y + row_idx * gap
        add_textbox(slide, 0.85, y - 0.02, 1.35, 0.35, label, font_size=12, bold=True, align=PP_ALIGN.RIGHT)
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(chart_x), Inches(y), Inches(chart_w), Inches(bar_h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        bg.line.color.rgb = COLOR_LIGHT_GRAY
        bar_w = chart_w * (float(value or 0) / max_value)
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(chart_x), Inches(y), Inches(bar_w), Inches(bar_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_ACCENT
        bar.line.color.rgb = COLOR_ACCENT
        add_textbox(slide, chart_x + chart_w + 0.15, y + 0.02, 1.2, 0.35, f"{value:g} {units}", font_size=13, bold=True)

    draw_bar(0, label1, value1)
    draw_bar(1, label2, value2)


def add_hyperlink_textbox(slide, x: float, y: float, w: float, h: float, text: str, url: str, font_size: int = 15):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = _safe_text(text)
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = COLOR_ACCENT
    if url:
        r.hyperlink.address = url
    return shape


def make_qr_image(url: str) -> BytesIO:
    qr = qrcode.QRCode(
        version=None,
        error_correction=qrcode.constants.ERROR_CORRECT_M,
        box_size=10,
        border=2,
    )
    qr.add_data(url or "")
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    bio = BytesIO()
    img.save(bio, format="PNG")
    bio.seek(0)
    return bio


def build_title_goal_slide(prs, deck):
    data = deck["title_goal"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.75, 0.6, 11.8, 0.55, data.get("session_title"), font_size=34, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, 1.28, 11.3, 0.85, data.get("article_title"), font_size=24, bold=True, align=PP_ALIGN.CENTER)
    add_small_label(slide, 1.25, 2.20, 10.9, "Teaching Goal")
    add_textbox(slide, 1.25, 2.45, 10.9, 0.95, data.get("teaching_goal"), font_size=20, align=PP_ALIGN.CENTER, fill=COLOR_ACCENT_LIGHT)
    add_section_label(slide, 2.65, 4.0, 8.0, "Five Questions Residents Should Answer")
    add_bullets(slide, 2.6, 4.55, 8.5, 1.55, _lines(data.get("five_questions")), font_size=17, bullet=False)
    add_footer(slide)
    return slide

def build_opening_case_slide(prs, deck):
    data = deck["opening_case"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Opening Patient Case")
    add_small_label(slide, 0.75, 0.98, 11.9, "Patient Case")
    add_textbox(slide, 0.75, 1.18, 11.9, 1.5, data.get("case_stem"), font_size=20, fill=COLOR_LIGHT_GRAY)
    add_inline_label(slide, 0.75, 2.72, 4.0, "Opening Question")
    add_textbox(slide, 0.75, 2.98, 11.9, 0.45, data.get("question"), font_size=22, bold=True, color=COLOR_ACCENT)
    add_section_label(slide, 1.0, 3.45, 3.4, "Answer Choices")
    add_bullets(slide, 1.0, 3.92, 6.4, 1.3, _lines(data.get("answer_choices")), font_size=18, bullet=False)
    add_small_label(slide, 0.85, 5.42, 11.6, "Facilitator Prompt")
    add_textbox(slide, 0.85, 5.65, 11.6, 0.75, data.get("facilitator_prompt"), font_size=18, bold=True, fill=COLOR_WARNING_LIGHT, align=PP_ALIGN.CENTER)
    add_footer(slide)
    return slide

def build_patient_problem_slide(prs, deck):
    data = deck["patient_problem"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "The Patient Problem")
    add_inline_label(slide, 0.75, 0.98, 4.0, "Problem Framing")
    add_textbox(slide, 0.75, 1.26, 11.9, 0.72, data.get("headline"), font_size=24, bold=True, color=COLOR_ACCENT)
    add_section_label(slide, 0.8, 2.18, 3.2, "Clinical Problem")
    add_bullets(slide, 0.95, 2.72, 11.5, 2.05, _lines(data.get("problem_bullets")), font_size=20)
    add_small_label(slide, 0.85, 5.20, 11.6, "Discussion Question")
    add_textbox(slide, 0.85, 5.45, 11.6, 0.75, data.get("discussion_question"), font_size=21, bold=True, fill=COLOR_ACCENT_LIGHT, align=PP_ALIGN.CENTER)
    add_footer(slide)
    return slide

def build_pico_slide(prs, deck):
    data = deck["pico"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "The Study Question", "PICO")
    labels = [("Patient/Problem", "patient"), ("Intervention", "intervention"), ("Comparison", "comparison"), ("Outcome", "outcome")]
    y = 1.35
    for label, key in labels:
        add_section_label(slide, 0.75, y, 2.2, label)
        add_textbox(slide, 3.1, y - 0.02, 9.55, 0.55, data.get(key), font_size=15)
        y += 0.85
    add_small_label(slide, 0.85, 4.72, 11.6, "Plain-Language Study Question")
    add_textbox(slide, 0.85, 4.95, 11.6, 0.75, data.get("plain_question"), font_size=20, bold=True, fill=COLOR_LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_inline_label(slide, 0.85, 5.86, 4.0, "Discussion Question")
    add_textbox(slide, 0.85, 6.10, 11.6, 0.45, data.get("discussion_question"), font_size=18, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(slide)
    return slide

def build_study_design_slide(prs, deck):
    data = deck["study_design"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "What They Did")
    add_inline_label(slide, 0.75, 0.96, 4.0, "Study Design")
    add_textbox(slide, 0.75, 1.24, 11.9, 0.50, data.get("design"), font_size=23, bold=True, color=COLOR_ACCENT)
    add_section_label(slide, 0.75, 1.95, 3.0, "What That Means")
    add_bullets(slide, 0.9, 2.42, 5.7, 2.2, _lines(data.get("design_bullets")), font_size=14)
    add_section_label(slide, 6.95, 1.95, 2.7, "Who Was Included")
    add_bullets(slide, 7.05, 2.42, 5.2, 1.25, _lines(data.get("included")), font_size=14)
    add_section_label(slide, 6.95, 3.9, 2.7, "Important Exclusions")
    add_bullets(slide, 7.05, 4.35, 5.2, 1.55, _lines(data.get("excluded")), font_size=13)
    add_small_label(slide, 0.85, 6.03, 11.6, "Discussion Question")
    add_textbox(slide, 0.85, 6.25, 11.6, 0.5, data.get("discussion_question"), font_size=17, bold=True, fill=COLOR_ACCENT_LIGHT, align=PP_ALIGN.CENTER)
    add_footer(slide)
    return slide

def build_main_result_slide(prs, deck):
    data = deck["main_result"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "What They Found")
    add_inline_label(slide, 0.75, 0.92, 4.0, "Main Result")
    add_textbox(slide, 0.75, 1.20, 11.9, 0.36, data.get("main_result"), font_size=23, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    visual_type = data.get("visual_type", "Results table")

    if visual_type == "Results table":
        add_results_table(slide, data.get("results_table", []), y=1.62)
    elif visual_type == "Big-number card":
        add_big_number_card(slide, data.get("big_number", ""), data.get("big_number_caption", ""))
        add_bullets(slide, 1.2, 4.55, 11.0, 0.8, _lines(data.get("key_results"))[:3], font_size=16)
    elif visual_type == "Simple bar chart":
        add_simple_bar_chart(
            slide,
            data.get("chart_title", ""),
            data.get("chart_group_1_label", "Group 1"),
            float(data.get("chart_group_1_value") or 0),
            data.get("chart_group_2_label", "Group 2"),
            float(data.get("chart_group_2_value") or 0),
            data.get("chart_units", ""),
        )
        add_bullets(slide, 1.2, 3.9, 11.0, 1.1, _lines(data.get("key_results"))[:3], font_size=16)
    else:
        add_bullets(slide, 1.0, 1.75, 11.4, 2.6, _lines(data.get("key_results")), font_size=20)

    add_small_label(slide, 0.85, 5.30, 11.6, "Plain-Language Result")
    add_textbox(slide, 0.85, 5.55, 11.6, 0.55, data.get("plain_result"), font_size=17, bold=True, fill=COLOR_LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_inline_label(slide, 0.85, 6.15, 4.0, "Discussion Question")
    add_textbox(slide, 0.85, 6.35, 11.6, 0.42, data.get("discussion_question"), font_size=15, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(slide)
    return slide

def build_clinical_bottom_line_slide(prs, deck):
    data = deck["clinical_bottom_line"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "What Should We Do?")

    # Top clinical interpretation.
    add_small_label(slide, 0.75, 0.94, 11.9, "Clinical Bottom Line")
    add_textbox(
        slide,
        0.75,
        1.20,
        11.9,
        0.75,
        data.get("bottom_line"),
        font_size=17,
        bold=True,
        fill=COLOR_ACCENT_LIGHT,
        align=PP_ALIGN.CENTER,
    )

    # Evidence-for and caution columns.
    add_section_label(slide, 0.9, 2.25, 2.6, "I Trust It Because")
    add_bullets(slide, 0.95, 2.75, 5.55, 1.55, _lines(data.get("trust_bullets")), font_size=15)
    add_section_label(slide, 6.85, 2.25, 2.7, "I Am Cautious Because")
    add_bullets(slide, 6.9, 2.75, 5.55, 1.9, _lines(data.get("caution_bullets")), font_size=14)

    # Bottom translation section.
    # These small labels make it clear that the first statement is for resident
    # clinical action and the second is a family-facing explanation.
    add_textbox(
        slide,
        0.85,
        4.88,
        11.6,
        0.24,
        "What I Would Do Clinically",
        font_size=11,
        bold=True,
        color=COLOR_ACCENT,
        align=PP_ALIGN.LEFT,
        margin=0.0,
    )
    add_textbox(
        slide,
        0.85,
        5.15,
        11.6,
        0.62,
        data.get("practice_statement"),
        font_size=15,
        bold=True,
        fill=COLOR_LIGHT_GRAY,
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        0.85,
        5.93,
        11.6,
        0.24,
        "How I Would Explain This To A Family",
        font_size=11,
        bold=True,
        color=COLOR_ACCENT,
        align=PP_ALIGN.LEFT,
        margin=0.0,
    )
    family_box = add_textbox(
        slide,
        0.85,
        6.20,
        11.6,
        0.58,
        data.get("family_explanation"),
        font_size=12,
        fill=COLOR_WARNING_LIGHT,
        align=PP_ALIGN.CENTER,
    )
    family_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_footer(slide)
    return slide


def build_paper_framework_slide(prs, deck):
    data = deck["paper_framework"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "PAPER framework discussion")
    boxes = [
        ("P", "Patient Problem", data.get("patient_problem_answer")),
        ("A", "Article Type", data.get("article_type_answer")),
        ("P", "Primary Question/Outcome", data.get("primary_question_answer")),
        ("E", "Evidence Quality", data.get("evidence_quality_answer")),
        ("R", "Real-World Use", data.get("real_world_answer")),
    ]
    coords = [(0.75, 1.2), (5.0, 1.2), (9.25, 1.2), (2.85, 4.05), (7.1, 4.05)]
    for (letter, title, body), (x, y) in zip(boxes, coords):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.35), Inches(2.25))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        shape.line.color.rgb = COLOR_ACCENT
        add_textbox(slide, x + 0.1, y + 0.08, 0.45, 0.4, letter, font_size=22, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.55, y + 0.12, 2.55, 0.35, title, font_size=12, bold=True)
        add_textbox(slide, x + 0.18, y + 0.58, 3.0, 1.45, body, font_size=11)
    add_footer(slide)
    return slide


def build_month_skill_slide(prs, deck):
    data = deck["month_skill"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, data.get("skill_title", "Month 1 focus skill"))
    add_section_label(slide, 0.9, 1.25, 3.6, "Five Questions")
    add_bullets(slide, 0.95, 1.75, 5.3, 2.3, _lines(data.get("reading_questions")), font_size=18, bullet=False)
    add_section_label(slide, 6.8, 1.25, 3.5, "Use This Paper As An Example")
    add_bullets(slide, 6.85, 1.75, 5.6, 2.3, _lines(data.get("this_paper_summary")), font_size=14, bullet=False)
    add_small_label(slide, 0.85, 5.15, 11.6, "Teaching Pearl")
    add_textbox(slide, 0.85, 5.4, 11.6, 0.9, data.get("teaching_pearl"), font_size=18, bold=True, fill=COLOR_ACCENT_LIGHT, align=PP_ALIGN.CENTER)
    add_footer(slide)
    return slide

def build_apply_back_slide(prs, deck):
    data = deck["apply_back"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Apply Back To The Patient")
    add_inline_label(slide, 0.85, 1.12, 4.5, "Return-To-Case Question")
    add_textbox(slide, 0.85, 1.35, 11.6, 0.8, data.get("return_question"), font_size=24, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_section_label(slide, 3.0, 2.7, 7.3, "Closing Vote")
    add_bullets(slide, 3.15, 3.22, 7.3, 1.4, _lines(data.get("vote_options")), font_size=20, bullet=False)
    add_small_label(slide, 0.85, 5.30, 11.6, "Facilitator Synthesis")
    add_textbox(slide, 0.85, 5.55, 11.6, 0.8, data.get("facilitator_synthesis"), font_size=17, bold=True, fill=COLOR_WARNING_LIGHT, align=PP_ALIGN.CENTER)
    add_footer(slide)
    return slide

def build_final_bottom_line_slide(prs, deck):
    data = deck["final_bottom_line"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Final Bottom Line")

    final_summary = data.get("final_summary")
    resident_take_home = data.get("resident_take_home")

    summary_x = 0.95
    summary_y = 1.32
    summary_w = 11.45
    summary_font_size = 20

    add_small_label(slide, summary_x, summary_y - 0.25, summary_w, "Final Summary")

    summary_h = estimate_textbox_height(
        final_summary,
        width_inches=summary_w,
        font_size=summary_font_size,
        min_height=0.85,
        max_height=2.05,
        padding=0.28,
    )

    summary_box = add_textbox(
        slide,
        summary_x,
        summary_y,
        summary_w,
        summary_h,
        final_summary,
        font_size=summary_font_size,
        bold=True,
        fill=COLOR_ACCENT_LIGHT,
        align=PP_ALIGN.CENTER,
        margin=0.06,
    )

    # Center text vertically inside the blue box
    summary_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Move take-home section based on the actual blue-box height
    takehome_label_y = summary_y + summary_h + 0.45
    takehome_text_y = takehome_label_y + 0.58

    add_section_label(
        slide,
        2.25,
        takehome_label_y,
        8.8,
        "Resident Take-Home Sentence",
    )

    takehome_h = estimate_textbox_height(
        resident_take_home,
        width_inches=10.65,
        font_size=22,
        min_height=0.65,
        max_height=1.35,
        padding=0.16,
    )

    takehome_box = add_textbox(
        slide,
        1.35,
        takehome_text_y,
        10.65,
        takehome_h,
        resident_take_home,
        font_size=22,
        bold=True,
        color=COLOR_ACCENT,
        align=PP_ALIGN.CENTER,
        margin=0.04,
    )

    takehome_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_footer(slide)
    return slide

def build_feedback_slide(prs, deck):
    """Final fixed feedback slide.

    This slide is not editable in the Streamlit app. The displayed website link
    and QR code destination are intentionally fixed for the pediatric residency
    journal club feedback tool.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.8, 0.75, 11.8, 0.7, THANK_YOU_TITLE, font_size=36, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.3, 1.65, 10.7, 0.75, THANK_YOU_MESSAGE, font_size=21, align=PP_ALIGN.CENTER)

    qr_card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.72), Inches(2.72), Inches(3.9), Inches(3.25))
    qr_card.fill.solid()
    qr_card.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    qr_card.line.color.rgb = COLOR_ACCENT_LIGHT

    qr_stream = make_qr_image(REDCAP_QR_URL)
    slide.shapes.add_picture(qr_stream, Inches(5.25), Inches(3.0), Inches(2.85), Inches(2.85))

    add_textbox(slide, 2.0, 6.1, 9.3, 0.35, FEEDBACK_INSTRUCTION, font_size=17, bold=True, align=PP_ALIGN.CENTER)
    add_hyperlink_textbox(slide, 1.2, 6.52, 10.9, 0.38, REDCAP_DISPLAY_URL, REDCAP_DISPLAY_URL, font_size=13)
    add_footer(slide, "Journal Club feedback")
    return slide



def build_custom_slide(prs, custom_slide: Dict[str, Any], index: int):
    """Create one optional user-added slide.

    Optional slides are intentionally simple so the core journal club structure
    stays standardized while still giving residents a safe place for extra
    context when it is truly needed.
    """
    title = _safe_text(custom_slide.get("title")) or f"Optional Slide {index}"
    body = _safe_text(custom_slide.get("body"))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, "Optional additional discussion slide")

    add_small_label(slide, 0.85, 1.18, 11.6, "Additional Teaching Point")
    body_box = add_textbox(
        slide,
        0.85,
        1.45,
        11.6,
        4.75,
        body,
        font_size=20,
        bold=False,
        fill=COLOR_LIGHT_GRAY,
        align=PP_ALIGN.LEFT,
        margin=0.14,
    )
    body_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(
        slide,
        1.2,
        6.35,
        10.9,
        0.32,
        "Use this slide only when the standard journal club structure needs one additional teaching point.",
        font_size=11,
        color=COLOR_MID,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, f"Optional slide {index} • Journal Club Builder")
    return slide

def build_facilitator_notes_slide(prs, deck):
    """Create an appendix-style slide for facilitator notes.

    python-pptx does not reliably expose full speaker-notes editing across versions,
    so notes are exported as an editable final slide instead of hidden PowerPoint notes.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Facilitator notes appendix")
    notes = []
    patient_note = deck.get("patient_problem", {}).get("speaker_note", "")
    if patient_note:
        notes.append(f"Slide 1 patient problem: {patient_note}")
    opening_prompt = deck.get("opening_case", {}).get("facilitator_prompt", "")
    if opening_prompt:
        notes.append(f"Opening case: {opening_prompt}")
    family = deck.get("clinical_bottom_line", {}).get("family_explanation", "")
    if family:
        notes.append(f"Family-facing explanation: {family}")

    custom_slides = deck.get(CUSTOM_SLIDES_KEY, [])
    if isinstance(custom_slides, list):
        optional_index = 0
        for core_slide in SLIDES:
            for custom_slide in custom_slides:
                if not isinstance(custom_slide, dict):
                    continue
                if _custom_slide_insert_after(custom_slide) != core_slide["id"]:
                    continue
                optional_index += 1
                reason = _safe_text(custom_slide.get("reason"))
                title = _safe_text(custom_slide.get("title")) or f"Optional Slide {optional_index}"
                if reason:
                    placement = _core_slide_label(core_slide["id"])
                    notes.append(f"Optional slide {optional_index} after {placement} ({title}): {reason}")

    add_bullets(slide, 0.85, 1.25, 11.6, 5.25, notes, font_size=15)
    add_footer(slide)
    return slide


def build_powerpoint(deck: Dict[str, Dict[str, Any]], include_facilitator_notes: bool = True) -> BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    core_builders = [
        ("title_goal", build_title_goal_slide),
        ("opening_case", build_opening_case_slide),
        ("patient_problem", build_patient_problem_slide),
        ("pico", build_pico_slide),
        ("study_design", build_study_design_slide),
        ("main_result", build_main_result_slide),
        ("clinical_bottom_line", build_clinical_bottom_line_slide),
        ("paper_framework", build_paper_framework_slide),
        ("month_skill", build_month_skill_slide),
        ("apply_back", build_apply_back_slide),
        ("final_bottom_line", build_final_bottom_line_slide),
    ]

    optional_index = 0
    for slide_id, builder in core_builders:
        builder(prs, deck)
        for custom_slide in _ordered_custom_slides_after(deck, slide_id):
            if _safe_text(custom_slide.get("body")):
                optional_index += 1
                build_custom_slide(prs, custom_slide, optional_index)

    if include_facilitator_notes:
        build_facilitator_notes_slide(prs, deck)

    # Always end with the feedback slide so the last visible slide has the
    # REDCap link and QR code.
    build_feedback_slide(prs, deck)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def slide_labels() -> List[str]:
    return [slide["label"] for slide in SLIDES]
